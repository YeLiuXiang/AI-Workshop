import json
import os
from io import BytesIO

from PIL import Image
from pptx import Presentation
from pptx.util import Inches


def generate_ppt(
    plan_file: str,
    slide_images: list[str],
    output_file: str,
) -> str:
    with open(plan_file, "r", encoding="utf-8") as file_handle:
        plan = json.load(file_handle)

    aspect_ratio = plan.get("aspect_ratio", "16:9")
    if aspect_ratio == "16:9":
        slide_width = Inches(13.333)
        slide_height = Inches(7.5)
    elif aspect_ratio == "4:3":
        slide_width = Inches(10)
        slide_height = Inches(7.5)
    else:
        slide_width = Inches(13.333)
        slide_height = Inches(7.5)

    presentation = Presentation()
    presentation.slide_width = slide_width
    presentation.slide_height = slide_height
    blank_layout = presentation.slide_layouts[6]
    slides_info = plan.get("slides", [])

    for index, image_path in enumerate(slide_images):
        if not os.path.exists(image_path):
            return f"Error: Slide image not found: {image_path}"

        slide = presentation.slides.add_slide(blank_layout)

        with Image.open(image_path) as image:
            if image.mode in ("RGBA", "P"):
                image = image.convert("RGB")

            image_width, image_height = image.size
            image_aspect = image_width / image_height
            slide_aspect = slide_width / slide_height
            slide_width_emu = int(slide_width)
            slide_height_emu = int(slide_height)

            if image_aspect > slide_aspect:
                new_width_emu = slide_width_emu
                new_height_emu = int(slide_width_emu / image_aspect)
                left = Inches(0)
                top = Inches((slide_height_emu - new_height_emu) / 914400)
            else:
                new_height_emu = slide_height_emu
                new_width_emu = int(slide_height_emu * image_aspect)
                left = Inches((slide_width_emu - new_width_emu) / 914400)
                top = Inches(0)

            image_bytes = BytesIO()
            image.save(image_bytes, format="JPEG", quality=95)
            image_bytes.seek(0)

            slide.shapes.add_picture(
                image_bytes,
                left,
                top,
                Inches(new_width_emu / 914400),
                Inches(new_height_emu / 914400),
            )

        if index < len(slides_info):
            slide_info = slides_info[index]
            notes = []
            if slide_info.get("title"):
                notes.append(f"Title: {slide_info['title']}")
            if slide_info.get("subtitle"):
                notes.append(f"Subtitle: {slide_info['subtitle']}")
            if slide_info.get("key_points"):
                notes.append("Key Points:")
                for point in slide_info["key_points"]:
                    notes.append(f"  - {point}")
            if slide_info.get("speaker_notes"):
                notes.append("Speaker Notes:")
                for note in slide_info["speaker_notes"]:
                    notes.append(f"  - {note}")

            if notes:
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                if text_frame is not None:
                    text_frame.text = "\n".join(notes)

    presentation.save(output_file)
    return f"Successfully generated presentation with {len(slide_images)} slides to {output_file}"


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="Generate PPTX from slide images")
    parser.add_argument("--plan-file", required=True)
    parser.add_argument("--slide-images", nargs="+", required=True)
    parser.add_argument("--output-file", required=True)
    arguments = parser.parse_args()

    try:
        print(generate_ppt(arguments.plan_file, arguments.slide_images, arguments.output_file))
    except Exception as error:
        print(f"Error while generating presentation: {error}")

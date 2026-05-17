import argparse
import json
import os
import random
from datetime import datetime, timezone
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt

try:
    from PIL import Image
except ImportError:
    Image = None


EMU_PER_INCH = 914400
PROJECT_ROOT = Path(__file__).resolve().parents[1]
DEFAULT_ACCENT_IMAGE_DIR = Path(
    os.getenv(
        "PPT_ACCENT_IMAGE_DIR",
        str(
            PROJECT_ROOT / "AI picture"
            if (PROJECT_ROOT / "AI picture").exists()
            else PROJECT_ROOT / "assets" / "accent-images"
        ),
    )
)
SUPPORTED_IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".webp"}


def load_json(path: str) -> dict:
    return json.loads(Path(path).read_text(encoding="utf-8"))


def hex_to_rgb(value: str) -> RGBColor:
    cleaned = value.strip().replace("#", "")
    return RGBColor.from_string(cleaned)


def inches(value: float) -> int:
    return Inches(value)


def apply_fill(shape, color: RGBColor) -> None:
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color
    shape.line.fill.background()


def apply_line(shape, color: RGBColor | None = None, width_pt: float = 0.8) -> None:
    line = shape.line
    if color is None:
        line.fill.background()
        return
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.width = Pt(width_pt)


def add_textbox(slide, left: float, top: float, width: float, height: float):
    return slide.shapes.add_textbox(inches(left), inches(top), inches(width), inches(height))


def set_frame_text(
    text_frame,
    paragraphs: list[str],
    font_name: str,
    font_size: float,
    color: RGBColor,
    bold: bool = False,
    level: int = 0,
    space_after: float = 4,
    align=PP_ALIGN.LEFT,
    vertical_anchor=MSO_VERTICAL_ANCHOR.TOP,
) -> None:
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = vertical_anchor

    for index, line in enumerate(paragraphs):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = line
        paragraph.level = level
        paragraph.alignment = align
        paragraph.space_after = Pt(space_after)
        run = paragraph.runs[0]
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold


def content_box(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_color: RGBColor,
    radius: float,
    border_color: RGBColor | None = None,
):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        inches(left),
        inches(top),
        inches(width),
        inches(height),
    )
    apply_fill(shape, fill_color)
    apply_line(shape, border_color, width_pt=0.8 if border_color else 0)
    try:
        if len(shape.adjustments) > 0:
            shape.adjustments[0] = radius
    except Exception:
        pass
    return shape


def add_panel_text(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    paragraphs: list[str],
    font_name: str,
    font_size: float,
    color: RGBColor,
    bold: bool = False,
    align=PP_ALIGN.CENTER,
    vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
    space_after: float = 6,
):
    box = add_textbox(slide, left, top, width, height)
    set_frame_text(
        box.text_frame,
        paragraphs,
        font_name,
        font_size,
        color,
        bold=bold,
        align=align,
        vertical_anchor=vertical_anchor,
        space_after=space_after,
    )
    return box


def add_section_title(slide, left: float, top: float, width: float, title: str, tokens: dict, color: RGBColor | None = None):
    return add_panel_text(
        slide,
        left,
        top,
        width,
        0.3,
        [title],
        tokens["title_font"],
        16.5,
        color or tokens["primary"],
        bold=True,
        align=PP_ALIGN.LEFT,
        vertical_anchor=MSO_VERTICAL_ANCHOR.TOP,
        space_after=0,
    )


def add_body_text(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    paragraphs: list[str],
    tokens: dict,
    font_size: float = 14.5,
    color: RGBColor | None = None,
    bold: bool = False,
    vertical_anchor=MSO_VERTICAL_ANCHOR.TOP,
    space_after: float = 8,
):
    return add_panel_text(
        slide,
        left,
        top,
        width,
        height,
        paragraphs,
        tokens["body_font"],
        font_size,
        color or tokens["text"],
        bold=bold,
        align=PP_ALIGN.LEFT,
        vertical_anchor=vertical_anchor,
        space_after=space_after,
    )


def strip_label_prefix(value: str) -> str:
    if "：" in value:
        return value.split("：", 1)[1].strip()
    if ":" in value:
        return value.split(":", 1)[1].strip()
    return value.strip()


def split_labeled_items(value: str) -> tuple[str, list[str]]:
    if "：" in value:
        label, content = value.split("：", 1)
    elif ":" in value:
        label, content = value.split(":", 1)
    else:
        return value.strip(), []

    items = [item.strip() for item in content.split("/") if item.strip()]
    return label.strip(), items


def split_title_body(value: str) -> tuple[str, str]:
    text = value.strip()
    if "：" in text:
        title, body = text.split("：", 1)
        return title.strip(), body.strip()
    if ":" in text:
        title, body = text.split(":", 1)
        return title.strip(), body.strip()
    return text, ""


def sanitize_cover_lines(lines: list[str]) -> list[str]:
    sanitized = []
    for line in lines:
        if "小组：" in line:
            parts = [part.strip() for part in line.split("|") if "小组：" not in part]
            cleaned = " | ".join([part for part in parts if part])
            if cleaned:
                sanitized.append(cleaned)
            continue
        sanitized.append(line)
    return sanitized


def build_cover_metadata(slide_data: dict) -> list[str]:
    fields = slide_data.get("fields", {})
    customer_line = fields.get("customer_line", "")
    group_line = fields.get("group_line", "")

    metadata = []
    if customer_line:
        metadata.append(customer_line)

    if group_line:
        cleaned = group_line.replace("小组：", "").strip()
        cleaned = cleaned.replace(" | ", " · ")
        cleaned = cleaned.replace("日期：", "日期 ")
        if cleaned:
            metadata.append(cleaned)

    return metadata[:2]


def build_cover_badges(slide_data: dict) -> list[str]:
    fields = slide_data.get("fields", {})
    badges = []

    customer_line = fields.get("customer_line", "").strip()
    if customer_line:
        badges.append(customer_line)

    group_line = fields.get("group_line", "").strip()
    if group_line:
        badges.append(group_line.replace(" | ", " · "))

    workshop_line = fields.get("workshop_line", "").strip()
    if not workshop_line:
        slots = slide_data.get("slots", {})
        workshop_line = slots.get("point_3", "").strip()
    if workshop_line:
        badges.append(workshop_line)

    return badges[:3]


def build_cover_highlight(slide_data: dict) -> str:
    fields = slide_data.get("fields", {})
    highlight = str(fields.get("cover_value_line", "")).strip()
    if highlight:
        return highlight

    notes = slide_data.get("speaker_notes", [])
    if notes:
        first_note = str(notes[0]).strip()
        if first_note:
            return first_note

    return str(slide_data.get("subtitle", "")).strip()


def build_cover_workshop_label(slide_data: dict) -> str:
    fields = slide_data.get("fields", {}) if isinstance(slide_data.get("fields"), dict) else {}
    workshop_line = str(fields.get("workshop_line", "")).strip()
    if not workshop_line:
        workshop_line = str(slide_data.get("eyebrow", "")).strip()
    label = strip_label_prefix(workshop_line) if workshop_line else ""
    return label.replace("AI Discovery Card Workshop", "AI Discovery Workshop")


def draw_arrow(slide, left: float, top: float, width: float, height: float, color: RGBColor) -> None:
    arrow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
        inches(left),
        inches(top),
        inches(width),
        inches(height),
    )
    apply_fill(arrow, color)
    apply_line(arrow, None)


def draw_down_arrow(slide, left: float, top: float, width: float, height: float, color: RGBColor) -> None:
    arrow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.DOWN_ARROW,
        inches(left),
        inches(top),
        inches(width),
        inches(height),
    )
    apply_fill(arrow, color)
    apply_line(arrow, None)


def render_flow_row(slide, top: float, label: str, items: list[str], tokens: dict, fill_color: RGBColor) -> None:
    label_box = content_box(slide, 0.95, top, 1.15, 0.72, fill_color, tokens["panel_radius"], None)
    add_panel_text(
        slide,
        1.1,
        top + 0.11,
        0.82,
        0.46,
        [label],
        tokens["body_font"],
        12,
        RGBColor(255, 255, 255),
        bold=True,
        space_after=0,
    )

    count = max(len(items), 1)
    available_width = 9.85
    gap = 0.22
    block_width = min(1.72, (available_width - (count - 1) * gap) / count)
    start_left = 2.3
    arrow_width = 0.26
    for index, item in enumerate(items):
        left = start_left + index * (block_width + gap)
        content_box(slide, left, top, block_width, 0.72, tokens["panel"], tokens["panel_radius"], tokens["panel_border"])
        add_body_text(
            slide,
            left + 0.14,
            top + 0.12,
            block_width - 0.28,
            0.42,
            [item],
            tokens,
            font_size=11.5,
            vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
            space_after=0,
        )
        if index < count - 1:
            draw_arrow(slide, left + block_width + 0.03, top + 0.2, arrow_width, 0.26, tokens["secondary"])


def theme_tokens(theme: dict) -> dict:
    colors = theme.get("colors", {})
    fonts = theme.get("fonts", {})
    layout = theme.get("layout", {})
    return {
        "primary": hex_to_rgb(colors.get("primary", "C62828")),
        "secondary": hex_to_rgb(colors.get("secondary", "E91E63")),
        "text": hex_to_rgb(colors.get("text", "333333")),
        "muted": hex_to_rgb(colors.get("muted", "777777")),
        "background": hex_to_rgb(colors.get("background", "FFFFFF")),
        "panel": hex_to_rgb(colors.get("panel", "F8F4F4")),
        "panel_alt": hex_to_rgb(colors.get("panel_alt", "F3EDED")),
        "panel_border": hex_to_rgb(colors.get("panel_border", "EADDDD")),
        "title_font": fonts.get("title", "Arial"),
        "body_font": fonts.get("body", "Arial"),
        "accent_font": fonts.get("accent", fonts.get("body", "Arial")),
        "panel_radius": float(layout.get("panel_radius", 0.1)),
    }


def get_accent_image_spec(slide_data: dict) -> dict | None:
    fields = slide_data.get("fields", {}) if isinstance(slide_data.get("fields"), dict) else {}
    spec = fields.get("accent_image")
    if not isinstance(spec, dict) or not spec.get("enabled"):
        return None
    return spec


def list_accent_image_candidates(image_dir: Path) -> list[Path]:
    if not image_dir.exists() or not image_dir.is_dir():
        return []
    return [
        path
        for path in image_dir.iterdir()
        if path.is_file() and path.suffix.lower() in SUPPORTED_IMAGE_EXTENSIONS
    ]


def crop_image_to_buffer(image_path: Path, width: float, height: float):
    if Image is None:
        return None
    try:
        with Image.open(image_path) as image:
            image_width, image_height = image.size
            if image_width <= 0 or image_height <= 0:
                return None

            box_aspect = width / height
            image_aspect = image_width / image_height
            if image_aspect > box_aspect:
                crop_width = int(image_height * box_aspect)
                crop_left = max((image_width - crop_width) // 2, 0)
                crop_box = (crop_left, 0, crop_left + crop_width, image_height)
            else:
                crop_height = int(image_width / box_aspect)
                crop_top = max((image_height - crop_height) // 2, 0)
                crop_box = (0, crop_top, image_width, crop_top + crop_height)

            cropped = image.crop(crop_box).convert("RGB")
            buffer = BytesIO()
            cropped.save(buffer, format="PNG")
            buffer.seek(0)
            return buffer
    except Exception:
        return None


def add_cropped_picture(slide, image_path: Path, left: float, top: float, width: float, height: float) -> bool:
    if Image is None:
        slide.shapes.add_picture(str(image_path), inches(left), inches(top), width=inches(width))
        return True

    buffer = crop_image_to_buffer(image_path, width, height)
    if buffer is None:
        return False
    slide.shapes.add_picture(
        buffer,
        inches(left),
        inches(top),
        width=inches(width),
        height=inches(height),
    )
    return True


def choose_accent_image(image_dir: Path) -> Path | None:
    candidates = list_accent_image_candidates(image_dir)
    if not candidates:
        return None
    return random.choice(candidates)


def add_accent_image(slide, slide_data: dict, tokens: dict) -> None:
    spec = get_accent_image_spec(slide_data)
    if not spec:
        return

    image_dir = Path(str(spec.get("directory") or DEFAULT_ACCENT_IMAGE_DIR))
    image_path = choose_accent_image(image_dir)
    if image_path is None:
        return

    left = float(spec.get("left", 9.9))
    top = float(spec.get("top", 2.0))
    width = float(spec.get("width", 2.1))
    height = float(spec.get("height", 4.6))

    add_cropped_picture(slide, image_path, left, top, width, height)


def set_slide_background(slide, color: RGBColor) -> None:
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = color


def add_header(slide, title: str, tokens: dict, subtitle: str = "") -> None:
    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        inches(0.6),
        inches(0.55),
        inches(0.12),
        inches(0.9),
    )
    apply_fill(accent, tokens["secondary"])

    title_box = add_textbox(slide, 0.9, 0.5, 8.8, 0.7)
    set_frame_text(title_box.text_frame, [title], tokens["title_font"], 24, tokens["primary"], bold=True, space_after=0)

    if subtitle:
        subtitle_box = add_textbox(slide, 0.92, 1.08, 10.0, 0.35)
        set_frame_text(subtitle_box.text_frame, [subtitle], tokens["body_font"], 10.5, tokens["muted"], space_after=0)

    divider = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        inches(0.9),
        inches(1.45),
        inches(11.5),
        inches(0.03),
    )
    apply_fill(divider, tokens["panel_alt"])


def render_cover(slide, slide_data: dict, tokens: dict) -> None:
    set_slide_background(slide, tokens["background"])

    cover_image = choose_accent_image(DEFAULT_ACCENT_IMAGE_DIR)
    if cover_image is not None and add_cropped_picture(slide, cover_image, 0, 0, 2.7, 7.5):
        image_overlay = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            inches(0),
            inches(0),
            inches(2.7),
            inches(7.5),
        )
        apply_fill(image_overlay, tokens["primary"])
        image_overlay.fill.transparency = 0.90
        image_overlay.line.fill.background()
    else:
        left_band = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            inches(0),
            inches(0),
            inches(2.7),
            inches(7.5),
        )
        apply_fill(left_band, tokens["primary"])

    top_accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        inches(9.55),
        inches(0.55),
        inches(2.1),
        inches(0.54),
    )
    apply_fill(top_accent, tokens["panel_alt"])
    apply_line(top_accent, None)

    content_box(slide, 3.0, 0.88, 8.6, 5.2, tokens["panel"], tokens["panel_radius"], tokens["panel_border"])

    workshop_label = build_cover_workshop_label(slide_data)
    if workshop_label:
        content_box(slide, 9.18, 0.98, 2.02, 0.42, tokens["background"], 0.06, tokens["panel_border"])
        add_panel_text(
            slide,
            9.34,
            1.08,
            1.7,
            0.16,
            [workshop_label],
            tokens["accent_font"],
            8.6,
            tokens["primary"],
            bold=True,
            align=PP_ALIGN.CENTER,
            vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
            space_after=0,
        )

    metadata = build_cover_metadata(slide_data)
    if metadata:
        add_panel_text(
            slide,
            3.38,
            1.12,
            7.5,
            0.28,
            [" · ".join(metadata)],
            tokens["body_font"],
            10.5,
            tokens["muted"],
            align=PP_ALIGN.LEFT,
            vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
            space_after=0,
        )

    add_panel_text(
        slide,
        3.38,
        1.52,
        7.45,
        1.32,
        [slide_data.get("title", "")],
        tokens["title_font"],
        30,
        tokens["primary"],
        bold=True,
        align=PP_ALIGN.LEFT,
        vertical_anchor=MSO_VERTICAL_ANCHOR.TOP,
        space_after=0,
    )

    subtitle = slide_data.get("subtitle", "")
    if subtitle:
        add_panel_text(
            slide,
            3.4,
            2.76,
            6.4,
            0.34,
            [subtitle],
            tokens["body_font"],
            11.5,
            tokens["muted"],
            align=PP_ALIGN.LEFT,
            vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
        )

    highlight = build_cover_highlight(slide_data)
    if highlight:
        content_box(slide, 3.38, 3.2, 7.54, 0.98, tokens["panel_alt"], tokens["panel_radius"], tokens["panel_border"])
        add_panel_text(
            slide,
            3.62,
            3.46,
            7.05,
            0.36,
            [highlight],
            tokens["body_font"],
            14,
            tokens["text"],
            align=PP_ALIGN.LEFT,
            vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
            space_after=0,
        )

    metadata_specs = [
        (3.38, 4.52, 4.2, "客户与行业"),
        (7.78, 4.52, 2.88, "小组与日期"),
    ]
    for index, badge_text in enumerate(build_cover_badges(slide_data)[:2]):
        left, top, width, label = metadata_specs[index]
        content_box(slide, left, top, width, 0.72, tokens["panel_alt"], tokens["panel_radius"], tokens["panel_border"])
        add_panel_text(
            slide,
            left + 0.18,
            top + 0.09,
            width - 0.36,
            0.14,
            [label],
            tokens["accent_font"],
            8.5,
            tokens["muted"],
            align=PP_ALIGN.LEFT,
            vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
            space_after=0,
        )
        add_panel_text(
            slide,
            left + 0.18,
            top + 0.27,
            width - 0.36,
            0.26,
            [badge_text],
            tokens["body_font"],
            10.0,
            tokens["text"],
            align=PP_ALIGN.LEFT,
            vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
            space_after=0,
        )


def render_three_panels(slide, slide_data: dict, tokens: dict) -> None:
    set_slide_background(slide, tokens["background"])
    add_header(slide, slide_data.get("title", ""), tokens, slide_data.get("subtitle", ""))

    boxes = [
        (0.9, 1.9, 3.6, 4.85),
        (4.87, 1.9, 3.6, 4.85),
        (8.84, 1.9, 3.6, 4.85),
    ]
    key_points = slide_data.get("key_points", ["", "", ""])
    for index, (left, top, width, height) in enumerate(boxes):
        content_box(
            slide,
            left,
            top,
            width,
            height,
            tokens["panel"] if index != 1 else tokens["panel_alt"],
            tokens["panel_radius"],
            tokens["panel_border"],
        )
        text = key_points[index] if index < len(key_points) else ""
        add_panel_text(
            slide,
            left + 0.35,
            top + 0.4,
            width - 0.7,
            height - 0.8,
            [text],
            tokens["body_font"],
            16 if index == 0 else 15,
            tokens["text"],
            bold=index == 0,
            align=PP_ALIGN.LEFT,
            space_after=8,
        )


def render_opportunity(slide, slide_data: dict, tokens: dict) -> None:
    set_slide_background(slide, tokens["background"])
    add_header(slide, slide_data.get("title", ""), tokens, slide_data.get("subtitle", ""))
    fields = slide_data.get("fields", {})
    accent_spec = get_accent_image_spec(slide_data)

    sections = [
        ("机会判断", strip_label_prefix(fields.get("opportunity_line", "")), tokens["panel"]),
        ("为什么是现在", strip_label_prefix(fields.get("why_now_line", "")), tokens["panel_alt"]),
        ("价值支撑", strip_label_prefix(fields.get("proof_line", "")), tokens["panel"]),
    ]
    if accent_spec:
        for index, (title, body, fill_color) in enumerate(sections):
            left, top, width, height = 0.9, 1.95 + index * 1.6, 8.65, 1.35
            content_box(slide, left, top, width, height, fill_color, tokens["panel_radius"], tokens["panel_border"])
            add_section_title(slide, left + 0.22, top + 0.18, width - 0.44, title, tokens)
            add_body_text(
                slide,
                left + 0.22,
                top + 0.52,
                width - 0.44,
                0.55,
                [body],
                tokens,
                font_size=13.4,
                vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
            )
        add_accent_image(slide, slide_data, tokens)
        return

    boxes = [
        (0.9, 1.95, 3.6, 4.8),
        (4.87, 1.95, 3.6, 4.8),
        (8.84, 1.95, 3.51, 4.8),
    ]
    for (title, body, fill_color), (left, top, width, height) in zip(sections, boxes):
        content_box(slide, left, top, width, height, fill_color, tokens["panel_radius"], tokens["panel_border"])
        add_section_title(slide, left + 0.25, top + 0.25, width - 0.5, title, tokens)
        add_body_text(
            slide,
            left + 0.25,
            top + 0.9,
            width - 0.5,
            height - 1.25,
            [body],
            tokens,
            font_size=14.5,
            vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
        )


def render_process_and_pain(slide, slide_data: dict, tokens: dict) -> None:
    set_slide_background(slide, tokens["background"])
    add_header(slide, slide_data.get("title", ""), tokens)
    fields = slide_data.get("fields", {})

    content_box(slide, 0.9, 1.9, 5.55, 4.8, tokens["panel"], tokens["panel_radius"], tokens["panel_border"])
    add_section_title(slide, 1.25, 2.15, 4.85, "现状流程", tokens)
    add_body_text(
        slide,
        1.2,
        2.7,
        4.95,
        3.45,
        fields.get("process_lines", []),
        tokens,
        14.5,
        space_after=8,
    )

    content_box(slide, 6.75, 1.9, 5.6, 3.3, tokens["panel_alt"], tokens["panel_radius"], tokens["panel_border"])
    add_section_title(slide, 7.1, 2.15, 4.9, "关键痛点", tokens)
    add_body_text(
        slide,
        7.0,
        2.68,
        5.0,
        1.95,
        fields.get("pain_lines", []),
        tokens,
        14.5,
        space_after=8,
    )

    content_box(slide, 6.75, 5.45, 5.6, 1.25, tokens["primary"], tokens["panel_radius"], None)
    add_body_text(
        slide,
        7.1,
        5.76,
        4.9,
        0.42,
        [fields.get("focus_line", "")],
        tokens,
        font_size=15.2,
        color=RGBColor(255, 255, 255),
        bold=True,
        vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
        space_after=0,
    )


def render_card_flow(slide, slide_data: dict, tokens: dict) -> None:
    set_slide_background(slide, tokens["background"])
    add_header(slide, slide_data.get("title", ""), tokens)
    fields = slide_data.get("fields", {})

    rows = [
        split_labeled_items(fields.get("perception_line", "")),
        split_labeled_items(fields.get("analysis_line", "")),
        split_labeled_items(fields.get("action_line", "")),
    ]
    content_box(slide, 0.9, 1.92, 11.45, 4.95, tokens["panel_alt"], tokens["panel_radius"], tokens["panel_border"])

    render_flow_row(slide, 2.18, rows[0][0] or "业务场景", rows[0][1], tokens, tokens["primary"])
    render_flow_row(slide, 3.62, rows[1][0] or "理解分析", rows[1][1], tokens, tokens["secondary"])
    render_flow_row(slide, 5.06, rows[2][0] or "生成交互 / 执行闭环", rows[2][1], tokens, tokens["primary"])


def render_flow(slide, slide_data: dict, tokens: dict) -> None:
    set_slide_background(slide, tokens["background"])
    add_header(slide, "目标业务闭环", tokens)
    fields = slide_data.get("fields", {})

    content_box(slide, 0.9, 1.9, 11.45, 0.8, tokens["primary"], tokens["panel_radius"], None)
    add_section_title(slide, 1.2, 2.06, 1.35, "业务触发", tokens, RGBColor(255, 255, 255))
    add_body_text(
        slide,
        2.3,
        2.05,
        8.9,
        0.42,
        [strip_label_prefix(fields.get("trigger_line", ""))],
        tokens,
        font_size=15.5,
        color=RGBColor(255, 255, 255),
        bold=True,
        vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
        space_after=0,
    )

    flow_lines = fields.get("flow_lines", [])
    content_box(slide, 0.9, 3.05, 11.45, 2.95, tokens["panel_alt"], tokens["panel_radius"], tokens["panel_border"])
    add_section_title(slide, 1.2, 3.25, 2.0, "核心流程", tokens)

    count = max(min(len(flow_lines), 5), 1)
    start_left = 1.15
    step_top = 3.72
    step_height = 1.68
    gap = 0.12
    arrow_width = 0.24
    total_width = 10.95
    step_width = (total_width - (count - 1) * (gap + arrow_width)) / count

    for index, line in enumerate(flow_lines[:5]):
        left = start_left + index * (step_width + gap + arrow_width)
        fill = tokens["panel"] if index % 2 == 0 else tokens["background"]
        content_box(slide, left, step_top, step_width, step_height, fill, tokens["panel_radius"], tokens["panel_border"])
        content_box(slide, left + 0.14, step_top + 0.12, 0.54, 0.34, tokens["secondary"], 0.05, None)
        add_panel_text(
            slide,
            left + 0.21,
            step_top + 0.18,
            0.4,
            0.1,
            [f"0{index + 1}" if index < 9 else str(index + 1)],
            tokens["accent_font"],
            9.5,
            RGBColor(255, 255, 255),
            bold=True,
            align=PP_ALIGN.CENTER,
            vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
            space_after=0,
        )
        add_body_text(
            slide,
            left + 0.16,
            step_top + 0.58,
            step_width - 0.28,
            0.92,
            [line],
            tokens,
            font_size=10.2,
            vertical_anchor=MSO_VERTICAL_ANCHOR.TOP,
            space_after=0,
        )
        if index < count - 1:
            arrow_left = left + step_width + gap
            draw_arrow(slide, arrow_left, step_top + 0.53, arrow_width, 0.22, tokens["secondary"])

    content_box(slide, 0.9, 6.12, 11.45, 0.72, tokens["secondary"], tokens["panel_radius"], None)
    add_section_title(slide, 1.2, 6.28, 1.6, "闭环结果", tokens, RGBColor(255, 255, 255))
    add_body_text(
        slide,
        2.45,
        6.25,
        8.95,
        0.22,
        [strip_label_prefix(fields.get("closure_line", ""))],
        tokens,
        font_size=13.6,
        color=RGBColor(255, 255, 255),
        bold=True,
        vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
        space_after=0,
    )


def render_architecture(slide, slide_data: dict, tokens: dict) -> None:
    set_slide_background(slide, tokens["background"])
    add_header(slide, slide_data.get("title", ""), tokens)
    fields = slide_data.get("fields", {})
    architecture_blocks = fields.get("architecture_blocks", []) or fields.get("layer_lines", [])
    support_lines = fields.get("support_lines", [])
    platform_chips = fields.get("platform_chips", [])
    platform_line = strip_label_prefix(fields.get("platform_line", ""))
    closure_line = fields.get("closure_line", "")

    content_box(slide, 0.9, 1.86, 11.45, 0.82, tokens["primary"], tokens["panel_radius"], None)
    add_body_text(
        slide,
        1.2,
        2.06,
        10.9,
        0.3,
        [closure_line or fields.get("architecture_title", "方案架构")],
        tokens,
        font_size=14.8,
        color=RGBColor(255, 255, 255),
        bold=True,
        vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
        space_after=0,
    )

    content_box(slide, 0.9, 2.95, 11.45, 4.18, tokens["panel_alt"], tokens["panel_radius"], tokens["panel_border"])
    add_section_title(slide, 1.18, 3.16, 6.8, fields.get("architecture_title", "方案架构"), tokens)

    cards = architecture_blocks[:4] or [fields.get("architecture_title", "方案架构")]
    step_colors = [tokens["primary"], tokens["secondary"], tokens["secondary"], tokens["primary"]]
    card_left = 1.18
    card_top = 3.62
    card_width = 2.46
    card_height = 2.44
    card_gap = 0.28
    for index, line in enumerate(cards):
        title, body = split_title_body(line)
        left = card_left + index * (card_width + card_gap)

        content_box(
            slide,
            left,
            card_top,
            card_width,
            card_height,
            tokens["panel"] if index % 2 else tokens["background"],
            tokens["panel_radius"],
            tokens["panel_border"],
        )
        content_box(slide, left, card_top, card_width, 0.48, step_colors[index % len(step_colors)], 0.1, None)
        add_panel_text(
            slide,
            left + 0.12,
            card_top + 0.12,
            card_width - 0.24,
            0.14,
            [f"0{index + 1}" if index < 9 else str(index + 1)],
            tokens["accent_font"],
            9.2,
            RGBColor(255, 255, 255),
            bold=True,
            align=PP_ALIGN.LEFT,
            vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
            space_after=0,
        )
        add_panel_text(
            slide,
            left + 0.16,
            card_top + 0.66,
            card_width - 0.32,
            0.28,
            [title],
            tokens["accent_font"],
            16.2,
            tokens["text"],
            bold=True,
            align=PP_ALIGN.LEFT,
            vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
            space_after=0,
        )
        add_body_text(
            slide,
            left + 0.16,
            card_top + 1.1,
            card_width - 0.32,
            0.96,
            [body or title],
            tokens,
            font_size=11.2,
            vertical_anchor=MSO_VERTICAL_ANCHOR.TOP,
            space_after=0,
        )
        if index < len(cards) - 1:
            draw_arrow(slide, left + card_width + 0.06, card_top + 1.08, 0.16, 0.2, tokens["secondary"])

    content_box(slide, 0.9, 7.28, 11.45, 0.6, tokens["panel"], tokens["panel_radius"], tokens["panel_border"])
    add_section_title(slide, 1.18, 7.43, 2.0, "落地能力与产品支撑", tokens)
    for index, line in enumerate(support_lines[:2]):
        add_body_text(
            slide,
            3.05,
            7.34 + index * 0.19,
            5.22,
            0.16,
            [line],
            tokens,
            font_size=10.4,
            vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
            space_after=0,
        )

    chips = platform_chips or [item.strip() for item in platform_line.replace("优先产品栈：", "").split("/") if item.strip()] or [platform_line or "待补充"]
    chip_left = 8.62
    chip_top = 7.36
    chip_gap = 0.1
    chip_width = 1.28
    for index, item in enumerate(chips[:3]):
        left = chip_left + index * (chip_width + chip_gap)
        content_box(slide, left, chip_top, chip_width, 0.34, tokens["panel_alt"] if index % 2 else tokens["background"], 0.08, tokens["panel_border"])
        add_panel_text(
            slide,
            left + 0.08,
            chip_top + 0.09,
            chip_width - 0.16,
            0.1,
            [item],
            tokens["body_font"],
            8.6,
            tokens["text"],
            align=PP_ALIGN.CENTER,
            vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
            space_after=0,
        )


def render_prototype(slide, slide_data: dict, tokens: dict) -> None:
    set_slide_background(slide, tokens["background"])
    add_header(slide, slide_data.get("title", ""), tokens)
    fields = slide_data.get("fields", {})
    top_labels = fields.get("top_labels") or ["原型名称", "系统形态", "价值落点"]
    accent_spec = get_accent_image_spec(slide_data)

    if accent_spec:
        content_box(slide, 0.9, 1.95, 8.65, 1.85, tokens["panel_alt"], tokens["panel_radius"], tokens["panel_border"])
        add_section_title(slide, 1.2, 2.16, 7.9, "原型摘要", tokens)
        add_body_text(
            slide,
            1.2,
            2.58,
            7.95,
            0.98,
            [
                strip_label_prefix(fields.get("prototype_name_line", "")),
                strip_label_prefix(fields.get("prototype_surface_line", "")),
                strip_label_prefix(fields.get("prototype_value_line", "")),
            ],
            tokens,
            font_size=12.2,
            space_after=6,
        )

        content_box(slide, 0.9, 4.12, 4.18, 2.55, tokens["panel"], tokens["panel_radius"], tokens["panel_border"])
        add_section_title(slide, 1.15, 4.37, 3.6, fields.get("flow_title", "核心流程"), tokens)
        add_body_text(slide, 1.15, 4.8, 3.6, 1.55, fields.get("prototype_flow_lines", []), tokens, 12.9, space_after=6)

        content_box(slide, 5.37, 4.12, 4.18, 2.55, tokens["panel_alt"], tokens["panel_radius"], tokens["panel_border"])
        add_section_title(slide, 5.62, 4.37, 3.6, fields.get("scope_title", "工作台范围"), tokens)
        add_body_text(slide, 5.62, 4.8, 3.6, 1.55, fields.get("prototype_scope_lines", []), tokens, 12.9, space_after=6)

        add_accent_image(slide, slide_data, tokens)
        return

    top_sections = [
        (top_labels[0], strip_label_prefix(fields.get("prototype_name_line", "")), tokens["panel"]),
        (top_labels[1], strip_label_prefix(fields.get("prototype_surface_line", "")), tokens["panel_alt"]),
        (top_labels[2], strip_label_prefix(fields.get("prototype_value_line", "")), tokens["panel"]),
    ]
    top_boxes = [
        (0.9, 1.95, 3.6, 2.2),
        (4.87, 1.95, 3.6, 2.2),
        (8.84, 1.95, 3.51, 2.2),
    ]
    for (title, body, fill_color), (left, top, width, height) in zip(top_sections, top_boxes):
        content_box(slide, left, top, width, height, fill_color, tokens["panel_radius"], tokens["panel_border"])
        add_section_title(slide, left + 0.25, top + 0.25, width - 0.5, title, tokens)
        add_body_text(
            slide,
            left + 0.25,
            top + 0.78,
            width - 0.5,
            height - 1.0,
            [body],
            tokens,
            font_size=13.6,
            vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
        )

    content_box(slide, 0.9, 4.48, 5.7, 2.2, tokens["panel_alt"], tokens["panel_radius"], tokens["panel_border"])
    add_section_title(slide, 1.15, 4.72, 5.1, fields.get("flow_title", "核心流程"), tokens)
    add_body_text(slide, 1.15, 5.22, 5.1, 1.12, fields.get("prototype_flow_lines", []), tokens, 13.7, space_after=6)

    content_box(slide, 6.65, 4.48, 5.7, 2.2, tokens["panel"], tokens["panel_radius"], tokens["panel_border"])
    add_section_title(slide, 6.9, 4.72, 5.1, fields.get("scope_title", "工作台范围"), tokens)
    add_body_text(slide, 6.9, 5.22, 5.1, 1.12, fields.get("prototype_scope_lines", []), tokens, 13.7, space_after=6)


def render_mapping(slide, slide_data: dict, tokens: dict) -> None:
    set_slide_background(slide, tokens["background"])
    add_header(slide, slide_data.get("title", ""), tokens)
    fields = slide_data.get("fields", {})

    content_box(slide, 0.9, 1.9, 7.45, 4.9, tokens["panel"], tokens["panel_radius"], tokens["panel_border"])
    add_section_title(slide, 1.2, 2.15, 6.85, "能力卡与产品映射", tokens)
    add_body_text(slide, 1.15, 2.72, 6.95, 3.45, fields.get("mapping_lines", []), tokens, 13.8, space_after=8)

    content_box(slide, 8.7, 1.9, 3.65, 2.1, tokens["panel_alt"], tokens["panel_radius"], tokens["panel_border"])
    add_section_title(slide, 9.0, 2.14, 3.05, "推荐产品栈", tokens)
    add_body_text(slide, 9.0, 2.58, 3.05, 1.02, [fields.get("platform_line", "")], tokens, 13.5, space_after=0)

    content_box(slide, 8.7, 4.35, 3.65, 2.45, tokens["primary"], tokens["panel_radius"], None)
    add_section_title(slide, 9.0, 4.62, 3.05, "现场交付物", tokens, RGBColor(255, 255, 255))
    add_body_text(slide, 9.0, 5.02, 3.05, 1.3, [fields.get("deliverable_line", "")], tokens, 13.5, RGBColor(255, 255, 255), space_after=0)


def render_value(slide, slide_data: dict, tokens: dict) -> None:
    set_slide_background(slide, tokens["background"])
    add_header(slide, slide_data.get("title", ""), tokens)
    fields = slide_data.get("fields", {})

    sections = [
        ("业务价值", fields.get("value_lines", []), tokens["panel"]),
        ("关键指标", fields.get("kpi_lines", []), tokens["panel_alt"]),
        ("POC 验证点", fields.get("poc_lines", []), tokens["panel"]),
    ]
    for index, (section_title, lines, fill_color) in enumerate(sections):
        left = 0.9 + index * 3.98
        content_box(slide, left, 1.95, 3.55, 4.9, fill_color, tokens["panel_radius"], tokens["panel_border"])
        add_section_title(slide, left + 0.22, 2.18, 3.1, section_title, tokens)
        add_body_text(slide, left + 0.22, 2.7, 3.1, 3.55, lines, tokens, 14, space_after=8)


def render_next_steps(slide, slide_data: dict, tokens: dict) -> None:
    set_slide_background(slide, tokens["background"])
    add_header(slide, slide_data.get("title", ""), tokens)
    fields = slide_data.get("fields", {})
    accent_spec = get_accent_image_spec(slide_data)

    content_box(slide, 0.9, 1.9, 11.45, 0.8, tokens["primary"], tokens["panel_radius"], None)
    add_section_title(slide, 1.2, 2.06, 1.35, fields.get("scope_title", "POC 范围"), tokens, RGBColor(255, 255, 255))
    add_body_text(
        slide,
        2.3,
        2.05,
        8.9,
        0.42,
        [strip_label_prefix(fields.get("scope_line", ""))],
        tokens,
        15.2,
        RGBColor(255, 255, 255),
        bold=True,
        vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
        space_after=0,
    )

    if accent_spec:
        content_box(slide, 0.9, 3.0, 4.25, 3.7, tokens["panel"], tokens["panel_radius"], tokens["panel_border"])
        add_section_title(slide, 1.2, 3.25, 3.65, fields.get("stakeholder_title", "关键角色"), tokens)
        add_body_text(slide, 1.2, 3.8, 3.65, 2.35, fields.get("stakeholder_lines", []), tokens, 13.5, space_after=8)

        content_box(slide, 5.4, 3.0, 4.25, 3.7, tokens["panel_alt"], tokens["panel_radius"], tokens["panel_border"])
        add_section_title(slide, 5.7, 3.25, 3.65, fields.get("next_action_title", "后续动作"), tokens)
        add_body_text(slide, 5.7, 3.8, 3.65, 2.35, fields.get("next_action_lines", []), tokens, 13.5, space_after=8)

        add_accent_image(slide, slide_data, tokens)
        return

    content_box(slide, 0.9, 3.0, 5.45, 3.7, tokens["panel"], tokens["panel_radius"], tokens["panel_border"])
    add_section_title(slide, 1.2, 3.25, 4.85, fields.get("stakeholder_title", "关键角色"), tokens)
    add_body_text(slide, 1.2, 3.8, 4.85, 2.35, fields.get("stakeholder_lines", []), tokens, 14.5, space_after=10)

    content_box(slide, 6.9, 3.0, 5.45, 3.7, tokens["panel_alt"], tokens["panel_radius"], tokens["panel_border"])
    add_section_title(slide, 7.2, 3.25, 4.85, fields.get("next_action_title", "后续动作"), tokens)
    add_body_text(slide, 7.2, 3.8, 4.85, 2.35, fields.get("next_action_lines", []), tokens, 14.5, space_after=10)


def render_slide(slide, slide_data: dict, tokens: dict) -> None:
    layout_key = slide_data.get("layout_key") or slide_data.get("slide_type") or ""
    if layout_key in {"cover-workshop", "cover-basic", "cover"}:
        render_cover(slide, slide_data, tokens)
        return
    if layout_key == "opportunity-summary":
        render_opportunity(slide, slide_data, tokens)
        return
    if layout_key == "as-is-pain-map":
        render_process_and_pain(slide, slide_data, tokens)
        return
    if layout_key == "capability-selection":
        render_card_flow(slide, slide_data, tokens)
        return
    if layout_key == "ai-flow-steps":
        render_flow(slide, slide_data, tokens)
        return
    if layout_key == "architecture-diagram":
        render_architecture(slide, slide_data, tokens)
        return
    if layout_key == "prototype-concept":
        render_prototype(slide, slide_data, tokens)
        return
    if layout_key == "product-mapping-table":
        render_mapping(slide, slide_data, tokens)
        return
    if layout_key == "business-value":
        render_value(slide, slide_data, tokens)
        return
    if layout_key == "poc-next-step":
        render_next_steps(slide, slide_data, tokens)
        return
    render_three_panels(slide, slide_data, tokens)


def set_notes(slide, slide_data: dict) -> None:
    notes = []
    if slide_data.get("title"):
        notes.append(f"Title: {slide_data['title']}")
    for line in slide_data.get("speaker_notes", []):
        notes.append(f"- {line}")
    if not notes:
        return
    notes_frame = slide.notes_slide.notes_text_frame
    if notes_frame is not None:
        notes_frame.text = "\n".join(notes)


def write_pptx(plan: dict, theme: dict, output_path: str) -> None:
    presentation = Presentation()
    aspect_ratio = plan.get("aspect_ratio", "16:9")
    if aspect_ratio == "4:3":
        presentation.slide_width = Inches(10)
        presentation.slide_height = Inches(7.5)
    else:
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)

    blank_layout = presentation.slide_layouts[6]
    tokens = theme_tokens(theme)

    for slide_data in plan.get("slides", []):
        slide = presentation.slides.add_slide(blank_layout)
        render_slide(slide, slide_data, tokens)
        set_notes(slide, slide_data)

    now = datetime.now(timezone.utc)
    core_properties = presentation.core_properties
    core_properties.created = now
    core_properties.modified = now
    core_properties.author = "AI Workshop Pipeline"
    core_properties.last_modified_by = "AI Workshop Pipeline"
    core_properties.title = str(plan.get("title") or Path(output_path).stem)
    core_properties.subject = str(plan.get("subtitle") or "Workshop generated presentation")

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)


def main() -> None:
    parser = argparse.ArgumentParser(description="Render a presentation-plan JSON into a simple branded PPTX")
    parser.add_argument("--plan", required=True, help="Path to presentation-plan.json")
    parser.add_argument("--theme-file", required=True, help="Path to theme JSON")
    parser.add_argument("--output", required=True, help="Path to output PPTX")
    args = parser.parse_args()

    plan = load_json(args.plan)
    theme = load_json(args.theme_file)
    write_pptx(plan, theme, args.output)
    print(args.output)


if __name__ == "__main__":
    main()
